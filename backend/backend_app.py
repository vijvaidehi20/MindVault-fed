import os
import uuid
import time
import threading
from flask import Flask, request, jsonify
from flask_cors import CORS
import fitz # PyMuPDF for PDF
from pptx import Presentation # python-pptx for PowerPoint
import google.generativeai as genai
from dotenv import load_dotenv
load_dotenv()

app = Flask(__name__)
CORS(app)

# A temporary folder to store uploaded files
TEMP_UPLOAD_FOLDER = 'temp_uploads'
if not os.path.exists(TEMP_UPLOAD_FOLDER):
    os.makedirs(TEMP_UPLOAD_FOLDER)

# Configure the Gemini API client
gemini_api_key = os.getenv("GEMINI_API_KEY")

if not gemini_api_key:
    # Raise an error if the key is not set, preventing a failed API call
    raise ValueError("GEMINI_API_KEY environment variable not set.")
genai.configure(api_key=gemini_api_key)

def extract_text_from_pdf(pdf_path):
    """Extracts text from a PDF file stream."""
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(pptx_path):
    """Extracts text from a PowerPoint file stream."""
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def summarize_text_with_gemini(text_content):
    """Sends text to the Gemini API for summarization."""
    if not text_content:
        return None
    
    # You can choose a different Gemini model if you want
    model = genai.GenerativeModel('gemini-1.5-pro')
    
    prompt = f"Summarize the following text concisely and clearly:\n\n{text_content}"
    
    try:
        response = model.generate_content(prompt)
        summary = response.text.strip()
        return summary
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
        return None

def delete_file_after_timeout(file_path, timeout=1800):  # 1800 seconds = 30 minutes
    """Deletes a file after a specified timeout."""
    def target():
        time.sleep(timeout)
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"Deleted temporary file: {file_path}")

    # Use a background thread to run the deletion task
    timer_thread = threading.Thread(target=target)
    timer_thread.daemon = True # Allows the thread to exit with the main program
    timer_thread.start()

@app.route('/api/upload', methods=['POST'])
def upload_file_endpoint():
    """Receives an uploaded file, saves it, and returns a unique ID."""
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    # Generate a unique ID to save the file
    unique_id = str(uuid.uuid4())
    file_extension = os.path.splitext(file.filename)[1].lower()
    file_path = os.path.join(TEMP_UPLOAD_FOLDER, f"{unique_id}{file_extension}")
    
    try:
        file.save(file_path)
        # Schedule the file for deletion
        delete_file_after_timeout(file_path)
        return jsonify({"message": "File uploaded successfully", "fileId": unique_id}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/summarize/<file_id>', methods=['GET'])
def summarize_endpoint(file_id):
    """Summarizes a file from the temporary storage using its unique ID."""
    try:
        # Find the file based on the unique ID
        matching_files = [f for f in os.listdir(TEMP_UPLOAD_FOLDER) if f.startswith(file_id)]
        if not matching_files:
            return jsonify({"error": "File not found or has expired"}), 404
        
        file_name = matching_files[0]
        file_path = os.path.join(TEMP_UPLOAD_FOLDER, file_name)
        file_extension = os.path.splitext(file_name)[1].lower()
        
        text_content = ""
        if file_extension == '.pdf':
            text_content = extract_text_from_pdf(file_path)
        elif file_extension == '.pptx' or file_extension == '.ppt':
            text_content = extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
        else:
            return jsonify({"error": "Unsupported file type"}), 415

        if not text_content:
            return jsonify({"error": "Could not extract text from file"}), 500

        summary = summarize_text_with_gemini(text_content)
        if not summary:
            return jsonify({"error": "Could not generate summary"}), 500
        
        return jsonify({"summary": summary})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(port=5000, debug=True)